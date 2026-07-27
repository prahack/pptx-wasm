#!/usr/bin/env python3
"""Generates the .pptx fixtures the golden suite renders.

One deck per milestone suite, each exercising exactly the features that milestone claims
to support. Fixtures are generated rather than committed so that a change to what a
milestone covers is a change to this file, reviewable alongside the code — and so the
repository does not accumulate binary blobs nobody can diff.

Run with the repo's venv:  ./.venv/bin/python fixtures/gen.py
"""

from __future__ import annotations

import io
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "generated"

# 16:9 at the PowerPoint 2013+ default size.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Layout indices in python-pptx's default template.
LAYOUT_TITLE = 0
LAYOUT_TITLE_CONTENT = 1
LAYOUT_TITLE_ONLY = 5
LAYOUT_BLANK = 6


def new_deck() -> Presentation:
    deck = Presentation()
    deck.slide_width = SLIDE_W
    deck.slide_height = SLIDE_H
    return deck


def blank(deck: Presentation):
    return deck.slides.add_slide(deck.slide_layouts[LAYOUT_BLANK])


def save(deck: Presentation, name: str) -> Path:
    path = OUT / f"{name}.pptx"
    deck.save(str(path))
    print(f"  {path.relative_to(ROOT.parent)}")
    return path


# --------------------------------------------------------------------------- m0

def gen_m0() -> None:
    """A blank slide. Proves the whole pipeline runs and the background is right."""
    deck = new_deck()
    blank(deck)
    save(deck, "m0-blank")


# --------------------------------------------------------------------------- m1

def gen_m1() -> None:
    """Two slides, one text box and one rectangle each — the thin vertical slice."""
    deck = new_deck()
    for i, (label, colour) in enumerate(
        [("First slide", RGBColor(0xC0, 0x30, 0x30)), ("Second slide", RGBColor(0x30, 0x60, 0xC0))]
    ):
        slide = blank(deck)

        box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(6.0), Inches(1.2))
        frame = box.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x20, 0x20, 0x20)

        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(4.0), Inches(2.0)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = colour
        rect.line.color.rgb = RGBColor(0x10, 0x10, 0x10)
        rect.line.width = Pt(2)
        # python-pptx puts an empty text frame on autoshapes; leave it empty so this
        # fixture stays about geometry.
        rect.text_frame.text = ""
        _ = i
    save(deck, "m1-basic")


# --------------------------------------------------------------------------- m2

def gen_m2() -> None:
    """Text: alignment, wrapping, weights, spacing, bullets, a non-system font."""
    deck = new_deck()

    slide = blank(deck)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(5.6), Inches(6.5))
    frame = box.text_frame
    frame.word_wrap = True

    specs = [
        ("Left aligned, and long enough that it has to wrap onto several lines.", PP_ALIGN.LEFT),
        ("Centred text that also wraps across more than one line of the box.", PP_ALIGN.CENTER),
        ("Right aligned, wrapping as well so the ragged edge is on the left.", PP_ALIGN.RIGHT),
        (
            "Justified text stretches the spaces so that both edges line up, which is a "
            "different code path from the other three.",
            PP_ALIGN.JUSTIFY,
        ),
    ]
    for i, (text, align) in enumerate(specs):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = text
        p.alignment = align
        p.font.size = Pt(16)
        p.space_after = Pt(10)

    # Weights and styles in one paragraph, so the line's metrics come from several fonts.
    mixed = frame.add_paragraph()
    mixed.space_before = Pt(12)
    for text, bold, italic, under in [
        ("Regular ", False, False, False),
        ("bold ", True, False, False),
        ("italic ", False, True, False),
        ("bold-italic ", True, True, False),
        ("underlined", False, False, True),
    ]:
        run = mixed.add_run()
        run.text = text
        run.font.size = Pt(18)
        run.font.bold = bold
        run.font.italic = italic
        run.font.underline = under

    # A font that is not installed on a typical Linux CI box, to exercise fallback.
    exotic = frame.add_paragraph()
    exotic.space_before = Pt(12)
    run = exotic.add_run()
    run.text = "Georgia at 20pt, to exercise the font fallback chain."
    run.font.name = "Georgia"
    run.font.size = Pt(20)

    # Bulleted and numbered content, via a content placeholder so the master's list
    # styles are in play — this is where inheritance meets text.
    slide2 = deck.slides.add_slide(deck.slide_layouts[LAYOUT_TITLE_CONTENT])
    slide2.shapes.title.text = "Bullets and numbering"
    body = slide2.placeholders[1].text_frame
    body.text = "First level bullet"
    for text, level in [
        ("Second level, which is indented further", 1),
        ("Third level, further still", 2),
        ("Back to the first level", 0),
        ("A line long enough that it must wrap inside its bullet's hanging indent", 0),
    ]:
        p = body.add_paragraph()
        p.text = text
        p.level = level

    # Line spacing and vertical anchoring.
    slide3 = blank(deck)
    for i, (spacing, anchor) in enumerate(
        [(1.0, MSO_ANCHOR.TOP), (1.5, MSO_ANCHOR.MIDDLE), (2.0, MSO_ANCHOR.BOTTOM)]
    ):
        box = slide3.shapes.add_textbox(
            Inches(0.5 + i * 4.2), Inches(0.5), Inches(4.0), Inches(6.0)
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = anchor
        p = frame.paragraphs[0]
        p.text = f"Line spacing {spacing}x, anchored {anchor}. Wrapping to several lines."
        p.line_spacing = spacing
        p.font.size = Pt(16)

    save(deck, "m2-text")


# --------------------------------------------------------------------------- m3

def _png(width: int, height: int) -> io.BytesIO:
    """A deterministic test image: coloured quadrants plus a diagonal.

    Quadrants make a crop obviously right or wrong; the diagonal makes rotation and flips
    obvious. A photograph would make every diff ambiguous.
    """
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (width, height), (255, 255, 255))
    d = ImageDraw.Draw(img)
    d.rectangle([0, 0, width // 2, height // 2], fill=(220, 60, 60))
    d.rectangle([width // 2, 0, width, height // 2], fill=(60, 160, 80))
    d.rectangle([0, height // 2, width // 2, height], fill=(60, 90, 200))
    d.rectangle([width // 2, height // 2, width, height], fill=(240, 190, 60))
    d.line([0, 0, width, height], fill=(20, 20, 20), width=max(2, width // 100))
    d.ellipse(
        [width * 0.35, height * 0.35, width * 0.65, height * 0.65],
        outline=(20, 20, 20),
        width=max(2, width // 120),
    )
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def gen_m3() -> None:
    """Preset geometry, connectors, groups with rotation and flips, images with crops."""
    deck = new_deck()

    # Preset geometry vocabulary.
    slide = blank(deck)
    shapes = [
        MSO_SHAPE.RECTANGLE,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MSO_SHAPE.OVAL,
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        MSO_SHAPE.RIGHT_TRIANGLE,
        MSO_SHAPE.DIAMOND,
        MSO_SHAPE.PENTAGON,
        MSO_SHAPE.HEXAGON,
        MSO_SHAPE.RIGHT_ARROW,
        MSO_SHAPE.LEFT_RIGHT_ARROW,
        MSO_SHAPE.CHEVRON,
        MSO_SHAPE.STAR_5_POINT,
        MSO_SHAPE.CROSS,
        MSO_SHAPE.CAN,
        MSO_SHAPE.CUBE,
        MSO_SHAPE.DONUT,
    ]
    palette = [
        RGBColor(0x44, 0x72, 0xC4),
        RGBColor(0xED, 0x7D, 0x31),
        RGBColor(0x70, 0xAD, 0x47),
        RGBColor(0xFF, 0xC0, 0x00),
    ]
    for i, shape_type in enumerate(shapes):
        col, row = i % 6, i // 6
        s = slide.shapes.add_shape(
            shape_type,
            Inches(0.4 + col * 2.1),
            Inches(0.4 + row * 2.2),
            Inches(1.8),
            Inches(1.8),
        )
        s.fill.solid()
        s.fill.fore_color.rgb = palette[i % len(palette)]
        s.line.color.rgb = RGBColor(0x20, 0x20, 0x20)
        s.line.width = Pt(1.5)
        s.text_frame.text = ""

    # Rotation and flips.
    slide2 = blank(deck)
    for i, rotation in enumerate([0, 30, 45, 90, 135, 180]):
        s = slide2.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(0.5 + i * 2.1),
            Inches(0.8),
            Inches(1.8),
            Inches(1.0),
        )
        s.rotation = rotation
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
        s.line.fill.background()
        s.text_frame.text = ""

    # Connectors between two shapes.
    a = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(3.5), Inches(2.0), Inches(1.2))
    b = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.0), Inches(5.2), Inches(2.0), Inches(1.2))
    for s in (a, b):
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0xE7, 0xE6, 0xE6)
        s.line.color.rgb = RGBColor(0x40, 0x40, 0x40)
        s.text_frame.text = ""
    from pptx.enum.shapes import MSO_CONNECTOR

    conn = slide2.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(3.0), Inches(4.1), Inches(8.0), Inches(5.8)
    )
    conn.line.color.rgb = RGBColor(0xC0, 0x30, 0x30)
    conn.line.width = Pt(3)

    # Images: whole, cropped, and stretched to a different aspect ratio.
    slide3 = blank(deck)
    slide3.shapes.add_picture(_png(400, 400), Inches(0.5), Inches(0.5), Inches(3.0), Inches(3.0))

    cropped = slide3.shapes.add_picture(
        _png(400, 400), Inches(4.2), Inches(0.5), Inches(3.0), Inches(3.0)
    )
    # Keep only the bottom-right quadrant, so a crop bug is unmissable.
    cropped.crop_left = 0.5
    cropped.crop_top = 0.5

    slide3.shapes.add_picture(_png(400, 200), Inches(0.5), Inches(4.0), Inches(6.7), Inches(1.6))

    save(deck, "m3-shapes")


# --------------------------------------------------------------------------- m4

def gen_m4() -> None:
    """Inheritance: placeholders that specify nothing and take everything from the chain.

    Every shape here is a placeholder used *as authored by the layout* — no position, no
    size, no font, no colour set on the slide. If the resolution chain is wrong, nothing
    lands in the right place, which is the point.
    """
    deck = new_deck()

    title_slide = deck.slides.add_slide(deck.slide_layouts[LAYOUT_TITLE])
    title_slide.shapes.title.text = "Quarterly Business Review"
    title_slide.placeholders[1].text = "Prepared by the Finance team"

    content = deck.slides.add_slide(deck.slide_layouts[LAYOUT_TITLE_CONTENT])
    content.shapes.title.text = "Highlights"
    body = content.placeholders[1].text_frame
    body.text = "Revenue up 12% year on year"
    for text, level in [
        ("Driven by the enterprise segment", 1),
        ("Renewals ahead of plan", 1),
        ("Costs held flat", 0),
        ("Headcount unchanged", 1),
    ]:
        p = body.add_paragraph()
        p.text = text
        p.level = level

    # A title-only layout, to exercise a different placeholder set on the same master.
    third = deck.slides.add_slide(deck.slide_layouts[LAYOUT_TITLE_ONLY])
    third.shapes.title.text = "Questions?"

    save(deck, "m4-template")


# --------------------------------------------------------------------------- m5a

def gen_m5a() -> None:
    """Tables: grid, header banding, merged cells, in-cell text."""
    deck = new_deck()
    slide = blank(deck)

    rows, cols = 5, 4
    table = slide.shapes.add_table(
        rows, cols, Inches(0.6), Inches(0.6), Inches(12.0), Inches(4.0)
    ).table

    headers = ["Region", "Q1", "Q2", "Q3"]
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.text_frame.paragraphs[0].font.bold = True

    data = [
        ["North", "1,204", "1,388", "1,502"],
        ["South", "988", "1,022", "1,140"],
        ["East", "1,455", "1,390", "1,610"],
        ["West", "720", "804", "868"],
    ]
    for r, row in enumerate(data, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text

    # A merged cell spanning the three quarter columns of the last row.
    table.cell(4, 1).merge(table.cell(4, 3))
    table.cell(4, 1).text = "Merged across three columns"

    # A second table with a vertical merge, since the two merge directions are separate
    # code paths.
    table2 = slide.shapes.add_table(3, 3, Inches(0.6), Inches(5.0), Inches(6.0), Inches(2.0)).table
    for r in range(3):
        for c in range(3):
            table2.cell(r, c).text = f"{r},{c}"
    table2.cell(0, 0).merge(table2.cell(2, 0))
    table2.cell(0, 0).text = "Vertical merge"

    save(deck, "m5a-tables")


# --------------------------------------------------------------------------- m5b

def gen_m5b() -> None:
    """Charts: clustered bar, line, and pie, with axes and legends.

    python-pptx writes the same `<c:numCache>`/`<c:strCache>` blocks PowerPoint does, so
    this exercises the code path a real deck takes — the cached values, not the embedded
    workbook.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    deck = new_deck()

    # Clustered column with two series.
    slide = blank(deck)
    data = CategoryChartData()
    data.categories = ["North", "South", "East", "West"]
    data.add_series("2023", (1204, 988, 1455, 720))
    data.add_series("2024", (1388, 1022, 1390, 804))
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(0.6), Inches(7.5), Inches(5.0), data
    )
    chart = frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.has_title = True
    chart.chart_title.text_frame.text = "Revenue by region"

    # Line chart on the same slide, so one golden covers both.
    line_data = CategoryChartData()
    line_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    line_data.add_series("Margin", (12.4, 13.1, 12.8, 14.6))
    line_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(8.4), Inches(0.6), Inches(4.4), Inches(2.4), line_data
    )
    line_frame.chart.has_legend = False

    # Pie, which has no axes and colours per point.
    pie_data = CategoryChartData()
    pie_data.categories = ["Enterprise", "Mid-market", "SMB"]
    pie_data.add_series("Mix", (0.55, 0.28, 0.17))
    pie_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(8.4), Inches(3.3), Inches(4.4), Inches(2.3), pie_data
    )
    pie_chart = pie_frame.chart
    pie_chart.has_legend = True
    pie_chart.legend.position = XL_LEGEND_POSITION.RIGHT
    pie_chart.legend.include_in_layout = False

    save(deck, "m5b-charts")


# --------------------------------------------------------------------------- m5c

def gen_m5c() -> None:
    """Effects: gradient fills, transparency, shadows."""
    deck = new_deck()
    slide = blank(deck)

    # python-pptx has no high-level gradient API for every case, but it does expose
    # `fill.gradient()`, which is enough for a linear two-stop gradient.
    grad = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.6), Inches(5.5), Inches(2.5)
    )
    grad.fill.gradient()
    grad.fill.gradient_angle = 45.0
    stops = grad.fill.gradient_stops
    stops[0].color.rgb = RGBColor(0x44, 0x72, 0xC4)
    stops[1].color.rgb = RGBColor(0xED, 0x7D, 0x31)
    grad.line.fill.background()
    grad.text_frame.text = ""

    # Transparency, via the XML directly: an alpha modifier on a solid fill.
    from pptx.oxml.ns import qn

    for i, alpha in enumerate([100_000, 60_000, 30_000]):
        s = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(7.0 + i * 1.6),
            Inches(0.8),
            Inches(2.2),
            Inches(2.2),
        )
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0xC0, 0x30, 0x30)
        s.line.fill.background()
        s.text_frame.text = ""
        srgb = s.fill.fore_color._xFill.find(qn("a:srgbClr"))
        if srgb is not None:
            from lxml import etree

            node = etree.SubElement(srgb, qn("a:alpha"))
            node.set("val", str(alpha))

    # A shape with an outer shadow.
    shadowed = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.0), Inches(4.0), Inches(2.0)
    )
    shadowed.fill.solid()
    shadowed.fill.fore_color.rgb = RGBColor(0x70, 0xAD, 0x47)
    shadowed.text_frame.text = ""
    sp_pr = shadowed._element.spPr
    from lxml import etree

    effect_lst = etree.SubElement(sp_pr, qn("a:effectLst"))
    shadow = etree.SubElement(effect_lst, qn("a:outerShdw"))
    shadow.set("blurRad", "76200")
    shadow.set("dist", "50800")
    shadow.set("dir", "2700000")
    shadow.set("rotWithShape", "0")
    colour = etree.SubElement(shadow, qn("a:srgbClr"))
    colour.set("val", "000000")
    alpha = etree.SubElement(colour, qn("a:alpha"))
    alpha.set("val", "40000")

    save(deck, "m5c-effects")


# --------------------------------------------------------------------------- m5d

def gen_m5d() -> None:
    """A gradient-filled title on a dark background.

    From a real bug: a gradient on *text* was being anchored to a box at the slide origin
    rather than to the glyphs, so every letter past the first inch sampled the clamped end
    stop and the whole title rendered one flat colour. Title slides built this way are
    common enough that it deserves a fixture.
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    deck = new_deck()
    slide = blank(deck)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x10, 0x25, 0x38)

    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(2.4))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "FIFA WORLD CUP AWARDS"
    run.font.size = Pt(54)
    run.font.bold = True

    # python-pptx has no gradient-on-text API, so the fill is written directly.
    rPr = run._r.get_or_add_rPr()
    grad = etree.SubElement(rPr, qn("a:gradFill"))
    gs_lst = etree.SubElement(grad, qn("a:gsLst"))
    for pos, colour in [(0, "4EA6DC"), (50000, "F2C744"), (100000, "4EDC94")]:
        gs = etree.SubElement(gs_lst, qn("a:gs"))
        gs.set("pos", str(pos))
        c = etree.SubElement(gs, qn("a:srgbClr"))
        c.set("val", colour)
    lin = etree.SubElement(grad, qn("a:lin"))
    lin.set("ang", "0")
    lin.set("scaled", "0")

    sub = slide.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(11.3), Inches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = "Celebrating Individual Excellence"
    sr.font.size = Pt(24)
    sr.font.bold = True
    sr.font.color.rgb = RGBColor(0x9F, 0xAF, 0xBF)

    save(deck, "m5d-gradient-text")


# --------------------------------------------------------------------------- m6

def gen_m6() -> None:
    """A large deck, for the performance bench. Not pixel-tested."""
    deck = new_deck()
    for i in range(60):
        slide = deck.slides.add_slide(deck.slide_layouts[LAYOUT_TITLE_CONTENT])
        slide.shapes.title.text = f"Slide {i + 1}"
        body = slide.placeholders[1].text_frame
        body.text = f"Point one on slide {i + 1}"
        for j in range(4):
            p = body.add_paragraph()
            p.text = f"Supporting detail {j + 1}, with enough words to require wrapping"
            p.level = 1 if j % 2 else 0
        for k in range(3):
            s = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(8.5),
                Inches(1.0 + k * 1.8),
                Inches(4.0),
                Inches(1.5),
            )
            s.fill.solid()
            s.fill.fore_color.rgb = RGBColor(0x44 + k * 0x30, 0x72, 0xC4)
            s.text_frame.text = f"Box {k + 1}"
    save(deck, "m6-large")


# --------------------------------------------------------------------------- bench

def gen_bench() -> None:
    """Two stress decks for `npm run bench`. Not pixel-tested.

    They probe the two axes that scale independently: how many slides a deck has, and how
    much is on one slide. The first is cheap by construction because parsing is lazy; the
    second is where the frame budget actually gets spent.
    """
    # Many slides, each ordinary.
    deck = new_deck()
    for i in range(250):
        slide = deck.slides.add_slide(deck.slide_layouts[LAYOUT_TITLE_CONTENT])
        slide.shapes.title.text = f"Slide {i + 1}"
        body = slide.placeholders[1].text_frame
        body.text = f"Point one on slide {i + 1}"
        for j in range(6):
            p = body.add_paragraph()
            p.text = f"Supporting detail {j + 1}, with enough words on the line to require wrapping"
            p.level = 1 if j % 2 else 0
        for k in range(4):
            s = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(8.5), Inches(0.6 + k * 1.6), Inches(4.0), Inches(1.4),
            )
            s.fill.solid()
            s.fill.fore_color.rgb = RGBColor(0x44 + k * 0x20, 0x72, 0xC4)
            s.text_frame.text = f"Box {k + 1}"
    save(deck, "bench-huge")

    # One slide, 2000 shapes: the density case that decides whether a frame is held.
    deck = new_deck()
    slide = blank(deck)
    for i in range(2000):
        col, row = i % 50, i // 50
        s = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.05 + col * 0.26), Inches(0.05 + row * 0.18),
            Inches(0.24), Inches(0.16),
        )
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor((i * 7) % 255, (i * 13) % 255, (i * 29) % 255)
        tf = s.text_frame
        tf.text = str(i)
        tf.paragraphs[0].font.size = Pt(6)
    save(deck, "bench-dense")


GENERATORS = {
    "m0": gen_m0,
    "m1": gen_m1,
    "m2": gen_m2,
    "m3": gen_m3,
    "m4": gen_m4,
    "m5a": gen_m5a,
    "m5b": gen_m5b,
    "m5c": gen_m5c,
    "m5d": gen_m5d,
    "m6": gen_m6,
    "bench": gen_bench,
}


def main(argv: list[str]) -> int:
    # The bench decks are large and slow to build, so they are opt-in.
    default = [k for k in GENERATORS if k != "bench"]
    wanted = [a for a in argv[1:] if not a.startswith("-")] or default
    unknown = [w for w in wanted if w not in GENERATORS]
    if unknown:
        print(f"unknown suite(s): {', '.join(unknown)}", file=sys.stderr)
        print(f"available: {', '.join(GENERATORS)}", file=sys.stderr)
        return 2

    if "--clean" in argv and OUT.exists():
        shutil.rmtree(OUT)
    OUT.mkdir(parents=True, exist_ok=True)

    print("Generating fixtures:")
    for name in wanted:
        GENERATORS[name]()
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
